# -*- coding: utf-8 -*-
import argparse
import csv
import os

from pptx import Presentation

from bible import TraditionalBibleContext, SimplifiedBibleContext


def parse_csv_file(_bc, book_name, book_title, chapter):
    csv_file_name = _bc.get_csv_file_name(book_name, chapter)
    pptx_file_name = _bc.get_pptx_file_name(book_name, chapter)

    if not _bc.to_be_updated(csv_file_name, pptx_file_name):
        # Skip
        return

    prs = Presentation(_bc.get_template_pptx())

    title_slide_layout_large = prs.slide_layouts[0]
    title_slide_layout_medium = prs.slide_layouts[1]
    title_slide_layout_small = prs.slide_layouts[2]
    title_slide_layout_extra_small = prs.slide_layouts[3]

    with open(csv_file_name) as csv_file:
        csv_reader = csv.reader(csv_file, delimiter=',')
        for row in csv_reader:
            verse_content = row[3]

            if len(verse_content) <= 150:  # 10 x 5
                slide = prs.slides.add_slide(title_slide_layout_large)
            elif len(verse_content) <= 216:  # 72 = 12 x 6
                slide = prs.slides.add_slide(title_slide_layout_medium)
            elif len(verse_content) <= 360:  # 72 = 15 x 8
                slide = prs.slides.add_slide(title_slide_layout_small)
            else:  # 18 x 10
                slide = prs.slides.add_slide(title_slide_layout_extra_small)

            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = u'{0} - [{1}:{2}]'.format(book_title, row[1], row[2])
            subtitle.text = row[3]
            # print "{0} {1}:{2} {3}".format(row[0], row[1], row[2], row[3])

    print "Creating {0}".format(pptx_file_name)
    prs.save(pptx_file_name)


def parse_csv_files(_bc, book_list):
    for book in book_list:
        for chapter in range(1, book['chapters'] + 1):
            parse_csv_file(_bc, book['name'], book['sk_title'] + " " + book['title'], chapter)


def convert(_bc):
    for template_file in os.listdir(_bc.template_path):
        if template_file.endswith(".pptx"):
            _bc.set_theme(template_file[0:-5])

            parse_csv_files(_bc, _bc.new_testiment_book_list)
            parse_csv_files(_bc, _bc.old_testiment_book_list)


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--theme", help="The pptx template theme. ")
    parser.add_argument("--force", help="Force update all PPTX files.", nargs='*')
    args = parser.parse_args()

    if args.theme is None:
        __all__ = True
    else:
        __all__ = False
        __theme__ = args.theme

    if args.force is None:
        __force_update__ = False
    else:
        __force_update__ = True

    if __all__:
        bible_context = TraditionalBibleContext()
        bible_context.force_update = __force_update__
        convert(bible_context)

        bible_context = SimplifiedBibleContext()
        bible_context.force_update = __force_update__
        convert(bible_context)

    else:
        pass
        # convert(BIBLE_CSV_PATH_TRADITIONAL, BIBLE_PPTX_PATH_TRADITIONAL, BIBLE_PPTX_TEMPLATE_PATH_TRADITIONAL)
