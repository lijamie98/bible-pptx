# -*- coding: utf-8 -*-
import argparse
import os
import sys

from pptx import Presentation

reload(sys)
sys.setdefaultencoding('utf8')

BIBLE_SCRIPT_PATH = os.path.dirname(os.path.realpath(__file__))
BIBLE_HOME = BIBLE_SCRIPT_PATH
BIBLE_DATABASE_PATH = os.path.abspath(os.path.join(BIBLE_HOME, "data", "text"))
BIBLE_PPTX_PATH = os.path.abspath(os.path.join(BIBLE_HOME, "data", "pptx_cut"))
BIBLE_PPTX_TEMPLATE_PATH = os.path.abspath(os.path.join(BIBLE_HOME, "data", "templates", "traditional"))

if not os.path.isdir(BIBLE_PPTX_PATH):
    os.makedirs(BIBLE_PPTX_PATH)

__theme__ = 'simple'

# Default force_update mode to False.
__force_update__ = False

max_verse = -1

book_dict = {
    "Samuel1": "1_Samuel",
    "Samuel2": "2_Samuel",
    "Kings1": "1_Kings",
    "Kings2": "2_Kings",
    "Chronicles1": "1_Chronicles",
    "Chronicles2": "2_Chronicles",
    "Corinthians1": "1_Corinthians1",
    "Corinthians2": "2_Corinthians1",
    "Thessalonians1": "2_Chronicles",
    "Thessalonians2": "2_Chronicles",
    "Timothy1": "1_Timothy",
    "Timothy2": "2_Timothy",
    "Peter1": "1_Peter",
    "Peter2": "2_Peter",
    "John1": "1_John",
    "John2": "2_John",
    "John3": "3_John"
}

last_book = ""
last_chapter = -1
prs = 0
title_slide_layout_large = 0
title_slide_layout_medium = 0
title_slide_layout_small = 0
title_slide_layout_extra_small = 0


def digest(book, book_chinese, chapter, verse, content):
    global last_book, last_chapter, prs, title_slide_layout_large, title_slide_layout_medium, title_slide_layout_small, title_slide_layout_extra_small
    if last_book != book or last_chapter != chapter:

        pptx_file_name = get_pptx_file_name(last_book, last_chapter)

        last_book = book
        last_chapter = chapter

        if not to_be_updated(pptx_file_name):
            # Skip
            return

        if prs != 0:
            print "Creating " + pptx_file_name + " ..."
            prs.save(pptx_file_name)

        prs = Presentation(os.path.join(BIBLE_PPTX_TEMPLATE_PATH, '{0}.pptx'.format(__theme__)))

        title_slide_layout_large = prs.slide_layouts[0]
        title_slide_layout_medium = prs.slide_layouts[1]
        title_slide_layout_small = prs.slide_layouts[2]
        title_slide_layout_extra_small = prs.slide_layouts[3]

    if len(content) <= 150:  # 10 x 5
        slide = prs.slides.add_slide(title_slide_layout_large)
    elif len(content) <= 216:  # 72 = 12 x 6
        slide = prs.slides.add_slide(title_slide_layout_medium)
    elif len(content) <= 360:  # 72 = 15 x 8
        slide = prs.slides.add_slide(title_slide_layout_small)
    else:  # 18 x 10
        slide = prs.slides.add_slide(title_slide_layout_extra_small)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = u'{0} - [{1}:{2}]'.format(book_chinese, chapter, verse)
    subtitle.text = content


def get_db_file_name():
    return os.path.join(BIBLE_DATABASE_PATH, 'cut.txt')


def get_pptx_file_name(book, chapter):
    return os.path.join(BIBLE_PPTX_PATH, __theme__, '{0}_{1:0>2d}.pptx'.format(book, chapter))


def to_be_updated(pptx_file_name):
    if __force_update__:
        return True

    if not os.path.exists(pptx_file_name):
        return True

    pptx_time = os.path.getmtime(pptx_file_name)
    db_time = os.path.getmtime(get_db_file_name())

    if db_time > pptx_time:
        return True

    return False


def convert():
    theme_path = os.path.join(BIBLE_PPTX_PATH, __theme__)

    if not os.path.isdir(theme_path):
        os.makedirs(theme_path)

    print get_db_file_name()

    with open(get_db_file_name()) as f:
        for line in f:
            words = line.split()
            if len(words) == 0:
                pass
            elif len(words) == 4:
                book = words[1]
                book_chinese = words[3]

                # Fix the book name
                if book in book_dict:
                    book = book_dict[book]
            elif len(words) == 5:
                ch_verse = words[3].split(":")
                content = words[4]
                chapter = int(ch_verse[0])
                verse = ch_verse[1]
                digest(book, book_chinese, chapter, verse, content)

    if prs != 0:
        pptx_file_name = get_pptx_file_name(last_book, last_chapter)
        prs.save(pptx_file_name)


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--theme", help="The pptx template theme. ")
    parser.add_argument("--force", help="Force update all PPTX files.", nargs='*')
    args = parser.parse_args()

    if args.theme is None:
        __all__ = True
    else:
        __all__ = False
        __theme__ = args.theme

    if args.force is None:
        __force_update__ = False
    else:
        __force_update__ = True

    if __all__:
        for theme_file in os.listdir(BIBLE_PPTX_TEMPLATE_PATH):
            if theme_file.endswith(".pptx"):
                __theme__ = theme_file[0:-5]
                convert()
    else:
        convert()
