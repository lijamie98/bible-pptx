# -*- coding: utf-8 -*-
import argparse
import csv
import os

from pptx import Presentation

from bible import BIBLE_CSV_PATH_TRADITIONAL, BIBLE_PPTX_PATH_TRADITIONAL, BIBLE_PPTX_TEMPLATE_PATH_TRADITIONAL, \
    BIBLE_DATABASE_PATH, BIBLE_COMBO_PPTX_PATH, \
    BIBLE_COMBO_PPTX_TEMPLATE_PATH
from bible.bible_book_list import new_testiment_book_list_traditional, old_testiment_book_list_traditional

__theme__ = 'simple'

# Default force_update mode to False.
__force_update__ = False

max_verse = -1


def get_combo_pptx_file_name(book, chapter):
    return os.path.join(BIBLE_COMBO_PPTX_PATH, __theme__, '{0}_{1:0>2d}.pptx'.format(book, chapter))


def parse_csv_file(book, chapter):
    csv_file_name = get_csv_file_name(book['name'], chapter)
    pptx_file_name = get_combo_pptx_file_name(book['name'], chapter)

    if not to_be_updated(csv_file_name, pptx_file_name):
        # Skip
        return

    prs = Presentation(os.path.join(BIBLE_COMBO_PPTX_TEMPLATE_PATH, '{0}.pptx'.format(__theme__)))

    title_slide_layout_large = prs.slide_layouts[0]
    title_slide_layout_medium = prs.slide_layouts[1]
    title_slide_layout_small = prs.slide_layouts[2]
    title_slide_layout_extra_small = prs.slide_layouts[3]

    with open(csv_file_name) as csv_file:
        csv_reader = csv.reader(csv_file, delimiter=',')
        for row in csv_reader:
            chapter = row[1];
            verse = row[2];
            verse_content = row[3]
            key = book['name'] + ":" + row[1] + ":" + row[2];
            if key in cu_bible:
                verse_cu = cu_bible[key]
            else:
                verse_cu = verse_content

            size = max(len(verse_content), len(verse_cu));

            if size <= 84:  # 6 x 7 * 2
                slide = prs.slides.add_slide(title_slide_layout_large)
            elif size <= 144:  # 8 x 9 * 2
                slide = prs.slides.add_slide(title_slide_layout_medium)
            elif size <= 220:  # 10 x 11 * 2
                slide = prs.slides.add_slide(title_slide_layout_small)
            else:  # 18 x 10
                slide = prs.slides.add_slide(title_slide_layout_extra_small)

            title = slide.shapes.title
            sk_verse_content = slide.placeholders[1]
            cu_verse_content = slide.placeholders[10]
            sk_book_title = slide.placeholders[11]
            cu_book_title = slide.placeholders[12]

            title.text = u'{0} [{1}:{2}]'.format(book['title'], chapter, verse)
            sk_verse_content.text = verse_content
            cu_verse_content.text = verse_cu
            sk_book_title.text = book['sk_title']
            cu_book_title.text = book['cu_title']
            # print "{0} {1}:{2} {3}".format(row[0], row[1], row[2], row[3])

    print "Creating {0}".format(pptx_file_name)
    prs.save(pptx_file_name)


def parse_csv_files(book_list):
    for book in book_list:
        for chapter in range(1, book['chapters'] + 1):
            parse_csv_file(book, chapter)


def convert():
    theme_path = os.path.join(BIBLE_COMBO_PPTX_PATH, __theme__)

    if not os.path.isdir(theme_path):
        os.makedirs(theme_path)

    parse_csv_files(new_testiment_book_list_traditional)
    parse_csv_files(old_testiment_book_list_traditional)


def get_db_file_name():
    return os.path.join(BIBLE_DATABASE_PATH, 'cut.txt')


book_dict = {
    "Samuel1": "1_Samuel",
    "Samuel2": "2_Samuel",
    "Kings1": "1_Kings",
    "Kings2": "2_Kings",
    "Chronicles1": "1_Chronicles",
    "Chronicles2": "2_Chronicles",
    "Corinthians1": "1_Corinthians",
    "Corinthians2": "2_Corinthians",
    "Thessalonians1": "1_Thessalonians",
    "Thessalonians2": "2_Thessalonians",
    "Timothy1": "1_Timothy",
    "Timothy2": "2_Timothy",
    "Peter1": "1_Peter",
    "Peter2": "2_Peter",
    "John1": "1_John",
    "John2": "2_John",
    "John3": "3_John",
    "Songs": "Song_of_Songs"
}

cu_bible = {}

if __name__ == "__main__":
    print "Loading CU database..."
    with open(get_db_file_name()) as f:
        for line in f:
            words = line.split()
            if len(words) == 0:
                pass
            elif len(words) == 4:
                book = words[1]
                book_chinese = words[3]

                # Fix the book name
                if book in book_dict:
                    book = book_dict[book]
            elif len(words) == 5:
                ch_verse = words[3].split(":")
                content = words[4]
                chapter = int(ch_verse[0])
                verse = ch_verse[1]
                cu_bible[book + ":" + str(chapter) + ":" + str(verse)] = content

    print "{} verses loaded".format(len(cu_bible))

    parser = argparse.ArgumentParser()
    parser.add_argument("--theme", help="The pptx template theme. ")
    parser.add_argument("--force", help="Force update all PPTX files.", nargs='*')
    args = parser.parse_args()

    if args.theme is None:
        __all__ = True
    else:
        __all__ = False
        __theme__ = args.theme

    if args.force is None:
        __force_update__ = False
    else:
        __force_update__ = True

    if __all__:
        for theme_file in os.listdir(BIBLE_COMBO_PPTX_TEMPLATE_PATH):
            if theme_file.endswith(".pptx"):
                __theme__ = theme_file[0:-5]
                convert()
    else:
        convert()
